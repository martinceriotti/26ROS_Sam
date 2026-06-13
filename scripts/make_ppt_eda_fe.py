"""
Genera la PPT de Feature Engineering del EDA para la presentación a la clase.
Run desde participant/:  python scripts/make_ppt_eda_fe.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Paleta California Republic ────────────────────────────────────────────────
RED    = RGBColor(0xBF, 0x0A, 0x30)   # carmesí California
BROWN  = RGBColor(0x8B, 0x65, 0x30)   # marrón del oso
GREEN  = RGBColor(0x2D, 0x7D, 0x32)   # verde del pasto
DARK   = RGBColor(0x3B, 0x23, 0x14)   # marrón oscuro (texto)
CREAM  = RGBColor(0xF5, 0xED, 0xE0)   # pergamino (fondo)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xEC, 0xE8, 0xE2)   # gris cálido para filas alternas

# ── Helpers ───────────────────────────────────────────────────────────────────
def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color, line_color=None, line_width=Pt(0)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_multiline(slide, lines, l, t, w, h, size=14, color=DARK, bold_first=False):
    """lines = list of (text, bold, color_override)"""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, bold, col = item, False, color
        else:
            text = item[0]
            bold = item[1] if len(item) > 1 else False
            col  = item[2] if len(item) > 2 else color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = col
    return txb

def table_cell(cell, text, bg, fg=DARK, size=11, bold=False, align=PP_ALIGN.LEFT):
    cell.fill.solid()
    cell.fill.fore_color.rgb = bg
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = fg

# ── Presentación ─────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]   # completamente en blanco


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Título
# ═══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
set_bg(s1, WHITE)

# Franja roja superior
add_rect(s1, 0, 0, 13.33, 0.22, RED)
# Franja roja inferior
add_rect(s1, 0, 7.28, 13.33, 0.22, RED)

# Bloque central con fondo pergamino
add_rect(s1, 0, 0.22, 13.33, 7.06, CREAM)

# Acento lateral izquierdo
add_rect(s1, 0, 0.22, 0.18, 7.06, RED)

# Título principal
add_text(s1, "Feature Engineering", 0.5, 1.6, 12.5, 1.2,
         size=48, bold=True, color=RED, align=PP_ALIGN.CENTER)
add_text(s1, "Exploración y nuevas variables en el EDA",
         0.5, 2.65, 12.5, 0.7, size=24, color=DARK, align=PP_ALIGN.CENTER)

# Línea separadora
add_rect(s1, 2.0, 3.5, 9.33, 0.05, BROWN)

# Subtítulos
add_text(s1, "Propiedades Residenciales — Miami / Sur de Florida",
         0.5, 3.7, 12.5, 0.5, size=16, color=BROWN, align=PP_ALIGN.CENTER)
add_text(s1, "Equipo SAM  ·  MDM Austral 2026  ·  Labo 2",
         0.5, 4.25, 12.5, 0.5, size=14, color=DARK, align=PP_ALIGN.CENTER, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE E1 — Enganche: "Adiviná el precio"
# ═══════════════════════════════════════════════════════════════════════════════
se1 = prs.slides.add_slide(BLANK)
set_bg(se1, CREAM)
add_rect(se1, 0, 0, 13.33, 0.22, RED)
add_rect(se1, 0, 7.28, 13.33, 0.22, RED)
add_rect(se1, 0, 0.22, 0.18, 7.06, RED)

add_text(se1, "¿CUÁNTO VALE ESTA CASA?",
         0.4, 0.45, 12.8, 0.85, size=34, bold=True, color=RED, align=PP_ALIGN.CENTER)

props = [
    ("📍", "Miami Beach, ZIP 33139"),
    ("🛏",  "4 habitaciones"),
    ("🚿",  "3 baños"),
    ("📐",  "2,100 ft²  (~195 m²)"),
    ("🏊",  "Piscina · Frente al agua"),
    ("🏠",  "Construida en 1987"),
]

for i, (icon, label) in enumerate(props):
    col = 0 if i < 3 else 6.8
    row = i % 3
    y = 1.5 + row * 1.2
    add_rect(se1, col + 0.4, y, 5.8, 1.05, WHITE, BROWN, Pt(1.5))
    add_text(se1, icon,  col + 0.55, y + 0.15, 0.8, 0.75, size=26, align=PP_ALIGN.CENTER)
    add_text(se1, label, col + 1.4,  y + 0.22, 4.6, 0.6,  size=18, color=DARK, bold=False)

add_rect(se1, 2.5, 6.5, 8.33, 0.65, RED)
add_text(se1, "💬  ¿Cuánto pagarías por esta propiedad?",
         2.5, 6.5, 8.33, 0.65, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE E2 — Enganche: meme Drake
# ═══════════════════════════════════════════════════════════════════════════════
se2 = prs.slides.add_slide(BLANK)
set_bg(se2, CREAM)
add_rect(se2, 0, 0, 13.33, 0.22, RED)
add_rect(se2, 0, 7.28, 13.33, 0.22, RED)
add_rect(se2, 0, 0.22, 0.18, 7.06, RED)

add_text(se2, "La forma tradicional vs. la nuestra",
         0.4, 0.32, 12.8, 0.6, size=20, bold=True, color=BROWN, align=PP_ALIGN.CENTER)

# Panel superior: NO
add_rect(se2, 1.2, 1.1, 10.8, 2.4, WHITE, RED, Pt(2))
add_rect(se2, 1.2, 1.1, 2.2, 2.4, RED)
add_text(se2, "✋\nNO", 1.2, 1.1, 2.2, 2.4, size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_multiline(se2, [
    ("Tasador humano", True, RED),
    ("• Cobra 3% de comisión sobre el precio de venta", False, DARK),
    ("• Tarda 2–5 días en entregar el informe", False, DARK),
    ("• Margen de error típico: ±15%", False, DARK),
], 3.7, 1.2, 7.8, 2.1, size=15)

# Panel inferior: SÍ
add_rect(se2, 1.2, 3.9, 10.8, 2.7, WHITE, GREEN, Pt(2))
add_rect(se2, 1.2, 3.9, 2.2, 2.7, GREEN)
add_text(se2, "👉\nSÍ", 1.2, 3.9, 2.2, 2.7, size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_multiline(se2, [
    ("Nuestro modelo ML", True, GREEN),
    ("• Costo: $0", False, DARK),
    ("• Tiempo: 3 segundos para 5,000 propiedades", False, DARK),
    ("• Entrenado con 12,000 transacciones reales de Miami", False, DARK),
    ("• Optimizado para ganar una simulación de inversión Monte Carlo", False, DARK),
], 3.7, 4.0, 7.8, 2.4, size=15)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE E3 — Enganche: stat de impacto
# ═══════════════════════════════════════════════════════════════════════════════
se3 = prs.slides.add_slide(BLANK)
set_bg(se3, CREAM)
add_rect(se3, 0, 0, 13.33, 0.22, RED)
add_rect(se3, 0, 7.28, 13.33, 0.22, RED)
add_rect(se3, 0, 0.22, 0.18, 7.06, RED)

add_text(se3, "$559,000",
         0.4, 1.0, 12.8, 2.2, size=96, bold=True, color=RED, align=PP_ALIGN.CENTER)

add_rect(se3, 2.0, 3.25, 9.33, 0.07, BROWN)

add_text(se3, "precio promedio de una casa en Miami en nuestro dataset.",
         0.5, 3.45, 12.3, 0.65, size=20, color=DARK, align=PP_ALIGN.CENTER)

add_text(se3, "¿Podemos predecirlo mejor que un agente inmobiliario?",
         0.5, 4.25, 12.3, 0.75, size=24, bold=True, color=BROWN, align=PP_ALIGN.CENTER)

add_rect(se3, 3.0, 5.3, 7.33, 0.65, RED)
add_text(se3, "Spoiler: sí.  Y lo hacemos gratis.",
         3.0, 5.3, 7.33, 0.65, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Calidad de datos: ceros como faltantes
# ═══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
set_bg(s2, WHITE)
add_rect(s2, 0, 0, 13.33, 0.22, RED)
add_rect(s2, 0, 7.28, 13.33, 0.22, RED)
add_rect(s2, 0, 0.22, 13.33, 7.06, CREAM)
add_rect(s2, 0, 0.22, 0.18, 7.06, RED)

# Header
add_rect(s2, 0.18, 0.22, 13.15, 0.85, RED)
add_text(s2, "01  ·  Calidad de Datos: ceros enmascarando faltantes",
         0.35, 0.28, 12.8, 0.72, size=22, bold=True, color=WHITE)

# Explicación
add_multiline(s2, [
    ("El dataset usa 0 en lugar de NaN cuando el dato no está disponible — patrón común en datasets de scraping inmobiliario.", False, DARK),
    ("Detectamos 3 columnas afectadas:", False, DARK),
], 0.4, 1.2, 12.5, 0.8, size=14)

# Tabla de ceros
headers = ["Columna", "NaN reales", "Ceros (missing real)", "Total missing", "Regla aplicada"]
rows = [
    ["bathrooms",  "225  (1.9%)", "75  (0.6%)", "300  (2.5%)", "Cero → NaN  (excepto LOT: terreno sin baños)"],
    ["livingArea", "172  (1.5%)", "7  (0.1%)",  "179  (1.5%)", "Cero → NaN  siempre (0 ft² es imposible)"],
    ["bedrooms",   "480  (4.1%)", "96  (0.8%)", "576  (4.9%)", "Cero → NaN  excepto CONDO/APARTMENT (studios válidos)"],
]

tbl = s2.shapes.add_table(4, 5, Inches(0.4), Inches(2.1), Inches(12.5), Inches(2.0)).table
col_widths = [1.4, 1.3, 1.7, 1.4, 6.7]
for i, w in enumerate(col_widths):
    tbl.columns[i].width = Inches(w)

for j, h in enumerate(headers):
    table_cell(tbl.cell(0, j), h, RED, WHITE, size=11, bold=True, align=PP_ALIGN.CENTER)

for i, row in enumerate(rows):
    bg = LGRAY if i % 2 == 0 else WHITE
    for j, val in enumerate(row):
        table_cell(tbl.cell(i+1, j), val, bg, DARK, size=10,
                   align=PP_ALIGN.CENTER if j < 4 else PP_ALIGN.LEFT)

# Solución
add_rect(s2, 0.4, 4.25, 12.5, 0.06, BROWN)
add_multiline(s2, [
    ("Solución:", True, RED),
    ("  Convertir ceros a NaN  →  imputar con mediana por homeType  →  fallback a mediana global.", False, DARK),
    ("  (Convención ya establecida en el proyecto para NaN reales.)", False, BROWN),
], 0.4, 4.35, 12.5, 1.0, size=13)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Los 7 features explorados
# ═══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
set_bg(s3, WHITE)
add_rect(s3, 0, 0, 13.33, 0.22, RED)
add_rect(s3, 0, 7.28, 13.33, 0.22, RED)
add_rect(s3, 0, 0.22, 13.33, 7.06, CREAM)
add_rect(s3, 0, 0.22, 0.18, 7.06, RED)

add_rect(s3, 0.18, 0.22, 13.15, 0.85, RED)
add_text(s3, "02  ·  7 Features Nuevos Explorados",
         0.35, 0.28, 12.8, 0.72, size=22, bold=True, color=WHITE)

# Tabla de features — 6 columnas con "¿Qué mide?"
headers3 = ["Feature", "Fórmula", "¿Qué mide?", "Corr.", "Leaky?", "Modelo"]
rows3 = [
    ["photo_signal",         "photoCount / mediana(photoCount, por tipo)",
     "Esfuerzo de marketing relativo: más fotos de lo habitual sugiere propiedad premium",
     "+0.152", "No", "✓ Sí"],
    ["hoa_burden",           "(hoa_fee_monthly × 12) / precio",
     "Cuota anual de mantenimiento como % del precio de compra",
     "−0.160", "Sí", "✗ No"],
    ["area_per_room",        "livingArea / (bedrooms + bathrooms)",
     "Espacio promedio por cuarto: diferencia studios pequeños de casas amplias",
     "+0.079", "No", "✓ Sí"],
    ["renovation_x_age",     "desc_mentions_renovated × property_age",
     "Captura el sobreprecio de propiedades viejas que fueron reformadas",
     "+0.075", "No", "✓ Sí"],
    ["listing_discount_pct", "(listing − precio) / listing",
     "% de descuento sobre el precio de lista al momento de cerrar la venta",
     "−0.071", "Sí", "✗ No"],
    ["market_premium",       "precio / taxAssessedValue",
     "Cuánto paga el mercado por encima del valor catastral oficial",
     "+0.019", "Sí", "✗ No"],
    ["price_per_sqft",       "precio / livingArea",
     "Precio por pie cuadrado — métrica estándar del sector inmobiliario",
     "+0.006", "Sí", "✗ No"],
]

tbl3 = s3.shapes.add_table(8, 6, Inches(0.4), Inches(1.2), Inches(12.5), Inches(5.5)).table
col_widths3 = [1.55, 2.5, 3.55, 0.9, 0.85, 0.9]
for i, w in enumerate(col_widths3):
    tbl3.columns[i].width = Inches(w)
tbl3.rows[0].height = Inches(0.45)

for j, h in enumerate(headers3):
    table_cell(tbl3.cell(0, j), h, RED, WHITE, size=11, bold=True, align=PP_ALIGN.CENTER)

for i, row in enumerate(rows3):
    bg = LGRAY if i % 2 == 0 else WHITE
    leaky = row[4] == "Sí"
    apto  = row[5].startswith("✓")
    for j, val in enumerate(row):
        fg = DARK
        if j == 4 and leaky:    fg = RED
        if j == 5 and apto:     fg = GREEN
        if j == 5 and not apto: fg = RED
        table_cell(tbl3.cell(i+1, j), val, bg, fg, size=9,
                   bold=(j in [0, 5]),
                   align=PP_ALIGN.CENTER if j in [3, 4, 5] else PP_ALIGN.LEFT)
    tbl3.rows[i+1].height = Inches(0.65)

add_text(s3,
         "* Leaky = el feature usa el precio de venta en su cálculo → se filtra el target al modelo.",
         0.4, 6.8, 12.5, 0.4, size=10, color=BROWN, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Los 3 features aptos + insights de mercado
# ═══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
set_bg(s4, WHITE)
add_rect(s4, 0, 0, 13.33, 0.22, RED)
add_rect(s4, 0, 7.28, 13.33, 0.22, RED)
add_rect(s4, 0, 0.22, 13.33, 7.06, CREAM)
add_rect(s4, 0, 0.22, 0.18, 7.06, RED)

add_rect(s4, 0.18, 0.22, 13.15, 0.85, RED)
add_text(s4, "03  ·  Features para el Modelo  +  Insights de Mercado",
         0.35, 0.28, 12.8, 0.72, size=22, bold=True, color=WHITE)

# Columna izquierda: features aptos
add_rect(s4, 0.4, 1.2, 6.2, 0.45, BROWN)
add_text(s4, "  Features aptos para Ronda 3",
         0.4, 1.2, 6.2, 0.45, size=14, bold=True, color=WHITE)

feat_cards = [
    ("photo_signal",     "+0.152", "¿Cuántas fotos tiene el anuncio comparado con casas similares?\nSi hay muchas más fotos de lo normal, suele ser una casa más cara."),
    ("area_per_room",    "+0.079", "¿Cuántos metros cuadrados le tocan a cada cuarto?\nNo es lo mismo un depto de 80m² con 2 cuartos que con 6."),
    ("renovation_x_age", "+0.075", "¿Es una casa vieja que fue refaccionada?\nUna casa antigua renovada vale más que una antigua sin tocar."),
]

for i, (name, corr, desc) in enumerate(feat_cards):
    y = 1.75 + i * 1.6
    add_rect(s4, 0.4, y, 6.2, 1.45, WHITE, RED, Pt(1.2))
    add_rect(s4, 0.4, y, 1.5, 1.45, RED)
    add_text(s4, corr, 0.4, y + 0.25, 1.5, 0.65, size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s4, name, 1.95, y + 0.08, 4.5, 0.35, size=13, bold=True, color=RED)
    add_text(s4, desc,  1.95, y + 0.45, 4.5, 0.95, size=10, color=DARK)

# Columna derecha: insights
add_rect(s4, 6.9, 1.2, 6.2, 0.45, BROWN)
add_text(s4, "  Insights del mercado (EDA)",
         6.9, 1.2, 6.2, 0.45, size=14, bold=True, color=WHITE)

insights = [
    ("1.81×", "Market Premium mediano",
     "El gobierno valúa las casas a $100 para cobrar impuestos,\npero en la calle se venden a $181. El catastro oficial\nqueda siempre muy por debajo del precio real."),
    ("56.4%", "propiedades vendidas sobre el precio de lista",
     "Más de la mitad de las casas se vendieron\nMÁS CARO que su precio publicado.\nEl precio del cartel es el mínimo, no el máximo."),
]

for i, (big, title, desc) in enumerate(insights):
    y = 1.75 + i * 1.75
    add_rect(s4, 6.9, y, 6.2, 1.6, WHITE, BROWN, Pt(1.2))
    add_rect(s4, 6.9, y, 1.5, 1.6, BROWN)
    add_text(s4, big,   6.9, y + 0.35, 1.5, 0.8, size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s4, title, 8.45, y + 0.08, 4.5, 0.38, size=12, bold=True, color=BROWN)
    add_text(s4, desc,  8.45, y + 0.5,  4.5, 1.0,  size=10, color=DARK)

# Nota al pie
add_rect(s4, 0.4, 6.9, 12.5, 0.06, RED)
add_text(s4,
         "Próximo paso: incorporar photo_signal, area_per_room y renovation_x_age en round2_segments.py → Ronda 3",
         0.4, 7.0, 12.5, 0.3, size=11, color=DARK, italic=True, align=PP_ALIGN.CENTER)


# ── Guardar ───────────────────────────────────────────────────────────────────
out = os.path.join("reports", "eda_feature_engineering.pptx")
prs.save(out)
print(f"PPT guardada: {out}")
