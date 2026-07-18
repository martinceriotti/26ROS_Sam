"""
Genera el PPT resumen — Equipo SAM
MDM Austral 2026 - Labo 2: Property Investment Competition

IMPORTANTE: esta presentacion se construye desde CERO (Presentation() en
blanco), no a partir de reports/UAustral-TextMiningClase01.pptx. Ese archivo
es de otra materia (Text Mining) y no se reutiliza su contenido ni el archivo
en si — solo tomamos como referencia visual el look and feel de Austral
(colores navy/naranja/rosa/verde, el logo oficial) que ya extrajimos como
imagenes propias en reports/assets/.

Contenido: resumen de la estrategia ganadora, que se probo, y que fracaso.

Run desde participant/:
    python scripts/generate_final_ppt.py
"""

import sys
from pathlib import Path
import warnings
warnings.filterwarnings('ignore')

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).parent))
import report_content as C

ROOT = Path(__file__).parent.parent
ASSETS = ROOT / "reports" / "assets"
OUT_PATH = ROOT / "reports" / "resumen_proyecto_SAM.pptx"

SLIDE_W, SLIDE_H = Emu(12192000), Emu(6858000)  # 13.333 x 7.5 in (16:9)

NAVY = RGBColor.from_string(C.NAVY.lstrip("#"))
ORANGE = RGBColor.from_string(C.ORANGE.lstrip("#"))
PINK = RGBColor.from_string(C.PINK.lstrip("#"))
TEAL = RGBColor.from_string(C.TEAL.lstrip("#"))
INDIGO = RGBColor.from_string(C.INDIGO.lstrip("#"))
GRAY = RGBColor.from_string(C.GRAY.lstrip("#"))
WHITE = RGBColor.from_string("FFFFFF")
DARK = RGBColor.from_string(C.DARK.lstrip("#"))
GREEN_OK = RGBColor.from_string(C.GREEN_OK.lstrip("#"))
RED_BAD = RGBColor.from_string(C.RED_BAD.lstrip("#"))

TRIANGLE_CYCLE = [PINK, TEAL, INDIGO]
VERDICT_RGB = {"great": GREEN_OK, "ok": INDIGO, "mixed": ORANGE, "bad": RED_BAD}
VERDICT_LABEL = {"great": "FUNCIONO", "ok": "OK", "mixed": "MIXTO", "bad": "NO FUNCIONO"}


# ─── Charts ─────────────────────────────────────────────────────────────────

def make_charts():
    plt.rcParams.update({"font.size": 12})

    fig, ax = plt.subplots(figsize=(9, 4.3))
    labels = [x[0] for x in C.ROUND9_COMPARISON]
    values = [x[1] for x in C.ROUND9_COMPARISON]
    colors = [C.GREEN_OK, C.RED_BAD, C.ORANGE]
    bars = ax.bar(labels, values, color=colors, alpha=0.9, edgecolor='black', linewidth=0.6)
    ax.set_ylabel("Mean ROI (%)")
    ax.set_title("Unico cambio: agregar features de texto (LLM)", color=C.NAVY, fontweight='bold')
    for b, v in zip(bars, values):
        ax.text(b.get_x() + b.get_width() / 2, v + 0.4, f"{v:.1f}%", ha='center', fontweight='bold')
    ax.spines[['top', 'right']].set_visible(False)
    ax.axhline(0, color='black', linewidth=0.6)
    fig.tight_layout()
    fig.savefig(ASSETS / "chart_round9.png", dpi=160, transparent=True)
    plt.close(fig)

    fig, ax = plt.subplots(figsize=(9.5, 4.6))
    scales = [x[0] for x in C.SCALE_SWEEP_LOCAL]
    rois = [x[1] for x in C.SCALE_SWEEP_LOCAL]
    ax.plot(scales, rois, '-o', color=C.NAVY, markersize=5, linewidth=2.2)
    ax.axvline(0.83, color=C.ORANGE, linestyle='--', linewidth=1.6, label='Escala elegida (0.83)')
    ax.axvline(1.00, color=C.GRAY, linestyle=':', linewidth=1.2, label='Escala original (1.00)')
    ax.set_xlabel("Factor de escala aplicado a la prediccion")
    ax.set_ylabel("Mean ROI (%) — simulacion local")
    ax.set_title("Barrido de escala sobre Ronda 8", color=C.NAVY, fontweight='bold')
    ax.legend(fontsize=10, loc='lower center')
    ax.spines[['top', 'right']].set_visible(False)
    ax.grid(alpha=0.25)
    fig.tight_layout()
    fig.savefig(ASSETS / "chart_scale_sweep.png", dpi=160, transparent=True)
    plt.close(fig)


# ─── Helpers de construccion (todo desde cero, sin depender de un template) ───

def add_triangle(slide, color, corner="tr"):
    if corner == "tr":
        shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Emu(9666111), Emu(-98779),
                                      Emu(2596444), Emu(2413000))
        shp.rotation = 180
    else:
        shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Emu(186158), Emu(-198852),
                                      Emu(1901656), Emu(2273969))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_logo(slide, corner="bl"):
    # Emu(1170432) de alto: el top tiene que dejar margen real antes del borde
    # inferior de la slide (Emu(6858000)) o el logo queda cortado.
    left = Emu(300000) if corner == "bl" else Emu(10145000)
    slide.shapes.add_picture(str(ASSETS / "austral_logo.jpeg"), left, Emu(5550000),
                              Emu(1746504), Emu(1170432))


def add_title(slide, text, size=30):
    tb = slide.shapes.add_textbox(Emu(700000), Emu(430000), Emu(10800000), Emu(650000))
    tf = tb.text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = NAVY
    r.font.name = "Calibri"
    return tb


def add_kicker(slide, text):
    tb = slide.shapes.add_textbox(Emu(700000), Emu(200000), Emu(6000000), Emu(280000))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = text
    r.font.size = Pt(13); r.font.color.rgb = ORANGE; r.font.bold = True
    r.font.name = "Calibri"


def add_pill(slide, verdict, top=Emu(1100000)):
    color = VERDICT_RGB[verdict]
    w = Emu(2100000)
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, SLIDE_W - Emu(300000) - w,
                                  top, w, Emu(400000))
    shp.fill.solid(); shp.fill.fore_color.rgb = WHITE
    shp.line.color.rgb = color; shp.line.width = Pt(1.5)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = VERDICT_LABEL[verdict]
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = color
    return shp


def add_bullets(slide, items, top=Emu(1700000), left=Emu(700000), width=Emu(10800000),
                 height=Emu(4500000), size=18, color=None, bullet_color=None):
    color = color or DARK
    bullet_color = bullet_color or ORANGE
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        p.line_spacing = 1.15
        r1 = p.add_run(); r1.text = "▸  "
        r1.font.size = Pt(size); r1.font.bold = True; r1.font.color.rgb = bullet_color
        r1.font.name = "Calibri"
        r2 = p.add_run(); r2.text = item
        r2.font.size = Pt(size); r2.font.color.rgb = color
        r2.font.name = "Calibri"
    return tb


def add_paragraphs(slide, paras, top=Emu(1700000), left=Emu(700000), width=Emu(10800000),
                    height=Emu(4500000), size=16, color=None):
    color = color or DARK
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        p.line_spacing = 1.2
        r = p.add_run(); r.text = para
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = "Calibri"
    return tb


def add_result_box(slide, text, color, top):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(700000), top,
                                  Emu(10800000), Emu(650000))
    box.fill.solid(); box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = color; box.line.width = Pt(1.5)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(180000)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = color
    return box


def base_slide(prs, tri_color, corner="tr"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_triangle(slide, tri_color, corner=corner)
    add_logo(slide, corner="bl" if corner == "tr" else "br")
    return slide


def set_cell_text(cell, text, size=11, bold=False, color=DARK, align=PP_ALIGN.LEFT,
                   fill=None):
    cell.margin_left = Emu(90000); cell.margin_right = Emu(90000)
    cell.margin_top = Emu(18000); cell.margin_bottom = Emu(18000)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Calibri"
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill


def main():
    make_charts()
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── 1. Portada ──────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
    bg.shadow.inherit = False
    for dx, dy, sz, col in [(9700000, -100000, 2600000, PINK),
                             (10900000, 1100000, 1600000, TEAL),
                             (8900000, 400000, 1000000, INDIGO)]:
        tri = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Emu(dx), Emu(dy), Emu(sz), Emu(sz))
        tri.rotation = 180
        tri.fill.solid(); tri.fill.fore_color.rgb = col; tri.line.fill.background()
        tri.shadow.inherit = False

    tb = s.shapes.add_textbox(Emu(700000), Emu(900000), Emu(8000000), Emu(500000))
    r = tb.text_frame.paragraphs[0].add_run(); r.text = "MDM Austral 2026 · Labo 2"
    r.font.size = Pt(16); r.font.color.rgb = ORANGE; r.font.bold = True

    tb2 = s.shapes.add_textbox(Emu(700000), Emu(1450000), Emu(9500000), Emu(1300000))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    r2 = tf2.paragraphs[0].add_run(); r2.text = C.PROJECT_TITLE
    r2.font.size = Pt(38); r2.font.bold = True; r2.font.color.rgb = WHITE

    tb3 = s.shapes.add_textbox(Emu(700000), Emu(2750000), Emu(9000000), Emu(500000))
    r3 = tb3.text_frame.paragraphs[0].add_run()
    r3.text = "Estrategia ganadora, pruebas y fracasos — " + C.TEAM
    r3.font.size = Pt(18); r3.font.color.rgb = RGBColor(0xD8, 0xDE, 0xEF)

    tb4 = s.shapes.add_textbox(Emu(700000), Emu(3250000), Emu(9000000), Emu(500000))
    r4 = tb4.text_frame.paragraphs[0].add_run()
    r4.text = f"Modelo final: {C.FINAL_MODEL}"
    r4.font.size = Pt(13); r4.font.color.rgb = RGBColor(0xB9, 0xC2, 0xDE); r4.font.italic = True

    # ── 2. Objetivo y mecanica ──────────────────────────────────────────────
    s = base_slide(prs, PINK, corner="tr")
    add_kicker(s, "EL JUEGO")
    add_title(s, "Objetivo y mecanica de la competencia")
    add_bullets(s, C.MECHANICS, size=16, top=Emu(1600000))

    # ── 3. Cronologia resumida: tabla con las 12 rondas + Kelly ─────────────
    s = base_slide(prs, TEAL, corner="bl")
    add_kicker(s, "RESUMEN")
    add_title(s, "Cronologia: que probamos, que funciono y que no")

    rows_data = [(r['n'], r['name'], r['verdict']) for r in C.ROUNDS]
    rows_data.append(("Extra", "Criterio de Kelly (escala por segmento) — casi empata al campeon", "ok"))
    rows_data.append(("Extra", "Edge por bucket (segmento x precio) — no supera al campeon", "mixed"))
    n_rows = len(rows_data) + 1
    row_h = Emu(350000)
    table_shape = s.shapes.add_table(n_rows, 3, Emu(700000), Emu(1500000),
                                      Emu(10800000), row_h * n_rows)
    table = table_shape.table
    table.columns[0].width = Emu(1600000)
    table.columns[1].width = Emu(7600000)
    table.columns[2].width = Emu(1600000)
    set_cell_text(table.cell(0, 0), "Ronda", size=12, bold=True, color=WHITE)
    set_cell_text(table.cell(0, 1), "Que se probo", size=12, bold=True, color=WHITE)
    set_cell_text(table.cell(0, 2), "Resultado", size=12, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER)
    for c in range(3):
        table.cell(0, c).fill.solid()
        table.cell(0, c).fill.fore_color.rgb = NAVY
    for i, (n, name, verdict) in enumerate(rows_data, start=1):
        fill = RGBColor(0xF7, 0xF8, 0xFC) if i % 2 == 0 else WHITE
        set_cell_text(table.cell(i, 0), n, size=10, color=GRAY, fill=fill)
        set_cell_text(table.cell(i, 1), name, size=10.5, color=DARK, fill=fill)
        set_cell_text(table.cell(i, 2), VERDICT_LABEL[verdict], size=9.5, bold=True,
                      color=VERDICT_RGB[verdict], align=PP_ALIGN.CENTER, fill=fill)
    for row in table.rows:
        row.height = row_h

    # ── 4. Estrategia ganadora ───────────────────────────────────────────────
    s = base_slide(prs, INDIGO, corner="tr")
    add_kicker(s, "ESTRATEGIA GANADORA")
    add_title(s, "Los 3 saltos que construyeron el modelo final")
    add_bullets(s, [
        "Embeddings de imagenes (CLIP): +92% de ROI vs el modelo tabular solo "
        "(el salto individual mas grande del proyecto).",
        "Regresion cuantil (percentil 35, no la media): ataca el error asimetrico "
        "del juego — sobreestimar destruye capital, subestimar solo cuesta la "
        "oportunidad.",
        "Fix quirurgico de propiedades distressed + calibracion final de escala "
        "(x0.83): confirmado en 3 corridas independientes de Practice, entre "
        "2x y 4x mejor Mean ROI que sin escalar.",
    ], size=17, top=Emu(1650000))
    add_result_box(s, f"Modelo final: {C.FINAL_MODEL}", NAVY, top=Emu(5600000))

    # ── 5. Lo que NO funciono: Ronda 9 texto/LLM ────────────────────────────
    s = base_slide(prs, RED_BAD, corner="bl")
    add_kicker(s, "LO QUE NO FUNCIONO")
    add_title(s, "Ronda 9: features de texto (LLM local)")
    s.shapes.add_picture(str(ASSETS / "chart_round9.png"), Emu(1450000), Emu(1550000),
                         width=Emu(9300000))
    add_paragraphs(s, [
        "Unico cambio respecto a Ronda 8: agregar 18 features semanticas extraidas "
        "de las descripciones con un LLM local (Ollama). Empeoro el modelo — y lo "
        "confirmamos con una comparacion controlada antes de descartarlo.",
    ], size=14, top=Emu(5900000), height=Emu(900000))

    # ── 6. La calibracion final ─────────────────────────────────────────────
    s = base_slide(prs, TEAL, corner="tr")
    add_kicker(s, "LA MEJORA FINAL")
    add_title(s, "Calibracion de escala sobre Ronda 8")
    s.shapes.add_picture(str(ASSETS / "chart_scale_sweep.png"), Emu(1350000), Emu(1500000),
                         width=Emu(9600000))
    add_paragraphs(s, [
        "Validado con 3 corridas independientes en el dashboard real: escala 0.83 "
        "le gano a Ronda 8 sin escalar por 2x a 4x en Mean ROI, sin excepciones.",
    ], size=14, top=Emu(6000000), height=Emu(800000))

    # ── 7. Criterio de Kelly ─────────────────────────────────────────────────
    s = base_slide(prs, INDIGO, corner="bl")
    add_kicker(s, "EXPERIMENTO NUEVO")
    add_title(s, "Criterio de Kelly, adaptado a la subasta")
    add_paragraphs(s, [
        "Kelly clasico asume perder = perder toda la apuesta. Acá no: si nos "
        "ganan la subasta, no perdemos nada. Lo adaptamos a p = Hit Rate, "
        "W = ganancia y L = perdida (como fraccion del costo), condicionado a "
        "ganar la subasta:  f* = p/L − q/W",
    ], size=15, top=Emu(1550000), height=Emu(1300000))

    headers = ["Segmento", "Hit Rate", "Kelly f*", "Escala sugerida"]
    tshape = s.shapes.add_table(4, 4, Emu(1400000), Emu(3050000), Emu(9400000), Emu(1600000))
    t = tshape.table
    for c, h in enumerate(headers):
        set_cell_text(t.cell(0, c), h, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        t.cell(0, c).fill.solid(); t.cell(0, c).fill.fore_color.rgb = NAVY
    for i, (seg, p, wgan, l, f, scale) in enumerate(C.KELLY_RESULTS, start=1):
        fill = RGBColor(0xF7, 0xF8, 0xFC) if i % 2 == 0 else WHITE
        vals = [seg, f"{p:.1f}%", f"{f:.2f}", f"{scale}"]
        for c, v in enumerate(vals):
            set_cell_text(t.cell(i, c), v, size=11.5, color=DARK, align=PP_ALIGN.CENTER, fill=fill)

    add_result_box(
        s, f"Local: 0.83 plano = {C.KELLY_FLAT_ROI:.2f}%  →  Kelly = {C.KELLY_SEGMENT_ROI:.2f}%  "
           f"·  Practice real: Kelly {C.PRACTICE_FINAL_VALIDATION[1][1]:.2f}% vs "
           f"campeon {C.PRACTICE_FINAL_VALIDATION[0][1]:.2f}% (empate practico)",
        TEAL, top=Emu(5000000),
    )
    add_paragraphs(s, [
        "Los 3 segmentos tienen perfil de riesgo parecido, asi que Kelly confirma "
        "que la escala plana ya estaba cerca del optimo — sin gran diferenciacion.",
    ], size=13, top=Emu(5750000), height=Emu(700000), color=GRAY)

    # ── 7b. Edge (valor esperado) por bucket fino ───────────────────────────
    s = base_slide(prs, TEAL, corner="tr")
    add_kicker(s, "EXPERIMENTO NUEVO")
    add_title(s, "Edge (valor esperado) por segmento x precio")
    add_paragraphs(s, [
        "EV = P(ganar) x (valor_real − costo pagado). Kelly ya midio esto por "
        "segmento y no encontro diferencia. Cortando mas fino — segmento x tercil "
        "de precio predicho — SI aparece un patron real y consistente.",
    ], size=15, top=Emu(1550000), height=Emu(1000000))

    headers2 = ["Bucket", "Hit Rate", "EV (% del bid)", "Escala"]
    tshape2 = s.shapes.add_table(4, 4, Emu(1600000), Emu(2650000), Emu(9000000), Emu(1500000))
    t2 = tshape2.table
    for c, h in enumerate(headers2):
        set_cell_text(t2.cell(0, c), h, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        t2.cell(0, c).fill.solid(); t2.cell(0, c).fill.fore_color.rgb = NAVY
    edge_summary = [
        ("bajo (los 3 segmentos)", 77.2, 1.51, "~0.80"),
        ("medio (los 3 segmentos)", 86.1, 1.75, "~0.83"),
        ("alto (los 3 segmentos)", 91.4, 2.06, "~0.87"),
    ]
    for i, (b, hit, ev, scale) in enumerate(edge_summary, start=1):
        fill = RGBColor(0xF7, 0xF8, 0xFC) if i % 2 == 0 else WHITE
        vals = [b, f"{hit:.1f}%", f"{ev:.2f}%", scale]
        for c, v in enumerate(vals):
            set_cell_text(t2.cell(i, c), v, size=11, color=DARK, align=PP_ALIGN.CENTER, fill=fill)

    m = C.PRACTICE_FINAL_VALIDATION
    add_result_box(
        s, f"Local: 0.83 plano = {C.EDGE_FLAT_ROI:.2f}%  →  Edge = {C.EDGE_BUCKET_ROI:.2f}%  "
           f"·  Practice real: Edge {m[2][1]:.2f}% vs campeon {m[0][1]:.2f}% (queda por debajo)",
        ORANGE, top=Emu(4550000),
    )
    add_paragraphs(s, [
        "El tercil 'alto' de precio siempre tiene mas edge que el 'bajo', en los 3 "
        "segmentos. Real, pero en Practice la mayor selectividad (menos volumen) "
        "no alcanzo para superar al campeon.",
    ], size=13, top=Emu(5300000), height=Emu(900000), color=GRAY)

    # ── 8. Recomendacion final ───────────────────────────────────────────────
    s = base_slide(prs, PINK, corner="tr")
    add_kicker(s, "RECOMENDACION FINAL")
    add_title(s, "Modelo elegido para la presentacion final")
    add_result_box(s, C.FINAL_MODEL, NAVY, top=Emu(1650000))
    add_paragraphs(s, [C.FINAL_RECOMMENDATION.split("\n\n")[1]],
                   top=Emu(2500000), size=15.5, height=Emu(1600000))
    add_bullets(s, [
        "Rondas 1 a 8: evolucion incremental, validada en Practice en cada paso.",
        "Ronda 9 (texto/LLM): pausada, no aporta con la cobertura/calidad actual.",
        "Kelly y Edge: probados a fondo, ninguno le gano al campeon — se cierra "
        "la busqueda de mejoras con la escala 0.83 confirmada.",
    ], top=Emu(4500000), size=14.5, height=Emu(1800000))

    # ── 9. Gracias ────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.shapes.add_picture(str(ASSETS / "austral_closing_bg.jpg"), 0, 0,
                         width=SLIDE_W, height=SLIDE_H)
    tb = s.shapes.add_textbox(Emu(700000), Emu(3000000), Emu(6000000), Emu(900000))
    r = tb.text_frame.paragraphs[0].add_run(); r.text = "MUCHAS GRACIAS"
    r.font.size = Pt(36); r.font.bold = True; r.font.color.rgb = NAVY

    prs.save(str(OUT_PATH))
    print(f"PPT generado: {OUT_PATH}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
